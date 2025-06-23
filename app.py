import os
import json
import zipfile
import mimetypes
import logging
from datetime import datetime
from pathlib import Path
from werkzeug.utils import secure_filename
from werkzeug.exceptions import RequestEntityTooLarge
from flask import Flask, request, redirect, url_for, render_template, flash, send_file, jsonify, abort
from flask_cors import CORS

# Document processing imports with proper error handling
try:
    import fitz  # PyMuPDF for PDF processing
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import openpyxl
    EXCEL_AVAILABLE = True
except ImportError:
    EXCEL_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from odf.opendocument import load
    from odf.text import P
    from odf.teletype import extractText
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False

# Configure logging
logging.basicConfig(level=logging.DEBUG)
logger = logging.getLogger(__name__)

app = Flask(__name__)
app.secret_key = os.environ.get("SESSION_SECRET", "default-secret-key-change-in-production")
CORS(app)

# Configuration
UPLOAD_FOLDER = 'static/uploads'
MAX_CONTENT_LENGTH = 10 * 1024 * 1024 * 1024  # 10GB limit
ALLOWED_EXTENSIONS = {
    'images': {'png', 'jpg', 'jpeg', 'gif', 'bmp', 'webp', 'svg', 'ico'},
    'documents': {'pdf', 'doc', 'docx', 'xls', 'xlsx', 'ppt', 'pptx', 'txt', 'rtf', 'odt', 'ods', 'odp'},
    'videos': {'mp4', 'avi', 'mkv', 'mov', 'wmv', 'flv', 'webm'},
    'audio': {'mp3', 'wav', 'ogg', 'aac', 'm4a', 'flac'},
    'archives': {'zip', 'rar', '7z', 'tar', 'gz', 'bz2'}
}

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

# Ensure upload directory exists
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

def get_file_type(filename):
    """Determine file type based on extension."""
    extension = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    for file_type, extensions in ALLOWED_EXTENSIONS.items():
        if extension in extensions:
            return file_type
    return 'other'

def allowed_file(filename):
    """Check if file extension is allowed."""
    if '.' not in filename:
        return False
    
    extension = filename.rsplit('.', 1)[1].lower()
    all_extensions = set()
    for extensions in ALLOWED_EXTENSIONS.values():
        all_extensions.update(extensions)
    
    return extension in all_extensions

def extract_text_from_file(file_path):
    """Extract text content from various file types."""
    try:
        file_extension = Path(file_path).suffix.lower()
        
        # Handle text files
        if file_extension == '.txt':
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='latin-1') as f:
                    return f.read()
        
        # Handle PDF files
        elif file_extension == '.pdf' and PDF_AVAILABLE:
            try:
                doc = fitz.open(file_path)
                text = ""
                for page in doc:
                    text += page.get_text()
                doc.close()
                return text
            except Exception as e:
                logger.error(f"Error processing PDF {file_path}: {str(e)}")
                return f"Error reading PDF: {str(e)}"
        
        # Handle Word documents
        elif file_extension in ['.doc', '.docx'] and DOCX_AVAILABLE:
    try:
        if file_extension == '.docx':
            doc = Document(file_path)
            text = []

            # Извлекаем текст из параграфов
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)

            # Извлекаем текст из таблиц
            for table in doc.tables:
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    if any(row_data):
                        text.append('	'.join(row_data))

            return '
'.join(text)
        else:
            return "Preview not available for .doc files. Please download to view."
    except Exception as e:
        logger.error(f"Error processing Word document {file_path}: {str(e)}")
        return f"Error reading Word document: {str(e)}"
        
        # Handle Excel files
        elif file_extension in ['.xls', '.xlsx'] and EXCEL_AVAILABLE:
            try:
                if file_extension == '.xlsx':
                    workbook = openpyxl.load_workbook(file_path)
                    text = []
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text.append(f"=== Sheet: {sheet_name} ===")
                        for row in sheet.iter_rows(values_only=True):
                            row_text = '\t'.join([str(cell) if cell is not None else '' for cell in row])
                            if row_text.strip():
                                text.append(row_text)
                    workbook.close()
                    return '\n'.join(text)
                else:
                    return "Preview not available for .xls files. Please download to view."
            except Exception as e:
                logger.error(f"Error processing Excel file {file_path}: {str(e)}")
                return f"Error reading Excel file: {str(e)}"
        
        # Handle PowerPoint files
        elif file_extension in ['.ppt', '.pptx'] and PPTX_AVAILABLE:
            try:
                if file_extension == '.pptx':
                    prs = Presentation(file_path)
                    text = []
                    for i, slide in enumerate(prs.slides, 1):
                        text.append(f"=== Slide {i} ===")
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text.append(shape.text)
                    return '\n'.join(text)
                else:
                    return "Preview not available for .ppt files. Please download to view."
            except Exception as e:
                logger.error(f"Error processing PowerPoint file {file_path}: {str(e)}")
                return f"Error reading PowerPoint file: {str(e)}"
        
        # Handle OpenDocument files
        elif file_extension in ['.odt', '.ods', '.odp'] and ODF_AVAILABLE:
            try:
                doc = load(file_path)
                text = []
                for paragraph in doc.getElementsByType(P):
                    text.append(extractText(paragraph))
                return '\n'.join(text)
            except Exception as e:
                logger.error(f"Error processing OpenDocument file {file_path}: {str(e)}")
                return f"Error reading OpenDocument file: {str(e)}"
        
        else:
            return "Preview not available for this file type. Please download to view."
    
    except Exception as e:
        logger.error(f"Unexpected error processing file {file_path}: {str(e)}")
        return f"Error processing file: {str(e)}"

def get_categories():
    """Get all categories and their structure."""
    categories = {}
    
    try:
        if not os.path.exists(UPLOAD_FOLDER):
            return categories
        
        for item in os.listdir(UPLOAD_FOLDER):
            item_path = os.path.join(UPLOAD_FOLDER, item)
            
            if os.path.isdir(item_path):
                categories[item] = {
                    'files': [],
                    'subcategories': {}
                }
                
                # Get files in category
                for file_item in os.listdir(item_path):
                    file_path = os.path.join(item_path, file_item)
                    
                    if os.path.isfile(file_path):
                        file_info = {
                            'name': file_item,
                            'type': get_file_type(file_item),
                            'size': os.path.getsize(file_path),
                            'modified': os.path.getmtime(file_path)
                        }
                        categories[item]['files'].append(file_info)
                    
                    elif os.path.isdir(file_path):
                        # Handle subcategories
                        subcategory_files = []
                        for sub_file in os.listdir(file_path):
                            sub_file_path = os.path.join(file_path, sub_file)
                            if os.path.isfile(sub_file_path):
                                file_info = {
                                    'name': sub_file,
                                    'type': get_file_type(sub_file), 
                                    'size': os.path.getsize(sub_file_path),
                                    'modified': os.path.getmtime(sub_file_path)
                                }
                                subcategory_files.append(file_info)
                        
                        categories[item]['subcategories'][file_item] = {
                            'files': subcategory_files
                        }
    
    except Exception as e:
        logger.error(f"Error getting categories: {str(e)}")
    
    return categories



def get_recent_files(limit=12):
    """Get recently uploaded files."""
    recent_files = []
    
    try:
        if not os.path.exists(UPLOAD_FOLDER):
            return recent_files
        
        all_files = []
        
        for root, dirs, files in os.walk(UPLOAD_FOLDER):
            for file in files:
                file_path = os.path.join(root, file)
                relative_path = os.path.relpath(file_path, UPLOAD_FOLDER)
                
                path_parts = relative_path.split(os.sep)
                category = path_parts[0] if len(path_parts) > 1 else 'general'
                subcategory = path_parts[1] if len(path_parts) > 2 else None
                
                file_info = {
                    'name': file,
                    'path': relative_path,
                    'category': category,
                    'subcategory': subcategory,
                    'type': get_file_type(file),
                    'size': os.path.getsize(file_path),
                    'modified': datetime.fromtimestamp(os.path.getmtime(file_path))
                }
                all_files.append(file_info)
        
        # Sort by modification time (newest first)
        all_files.sort(key=lambda x: x['modified'], reverse=True)
        recent_files = all_files[:limit]
    
    except Exception as e:
        logger.error(f"Error getting recent files: {str(e)}")

    
    return recent_files

@app.errorhandler(RequestEntityTooLarge)
def handle_file_too_large(e):
    flash('File too large. Maximum size is 10GB.', 'error')
    return redirect(url_for('index'))

@app.route('/')
def index():
    """Main page with file upload and category overview."""
    try:
        categories = get_categories()
        recent_files = get_recent_files()
        
        return render_template('index.html', 
                             categories=categories,
                             recent_files=recent_files)
    except Exception as e:
        logger.error(f"Error in index route: {str(e)}")
        flash(f'Error loading page: {str(e)}', 'error')
        return render_template('index.html', categories={}, recent_files=[])

@app.route('/upload', methods=['POST'])
def upload_files():
    """Handle file uploads."""
    try:
        if 'files' not in request.files:
            flash('No files selected', 'error')
            return redirect(url_for('index'))
        
        files = request.files.getlist('files')
        category = request.form.get('category', 'general').strip()
        subcategory = request.form.get('subcategory', '').strip()
        
        if not category:
            category = 'general'
        
        # Sanitize category and subcategory names
        category = secure_filename(category)
        if subcategory:
            subcategory = secure_filename(subcategory)
        
        uploaded_count = 0
        errors = []
        
        for file in files:
            if file.filename == '':
                continue
            
            try:
                if file and allowed_file(file.filename):
                    filename = secure_filename(file.filename)
                    
                    # Create directory structure
                    if subcategory:
                        upload_path = os.path.join(UPLOAD_FOLDER, category, subcategory)
                    else:
                        upload_path = os.path.join(UPLOAD_FOLDER, category)
                    
                    os.makedirs(upload_path, exist_ok=True)
                    
                    # Handle duplicate filenames
                    file_path = os.path.join(upload_path, filename)
                    counter = 1
                    base_name, extension = os.path.splitext(filename)
                    
                    while os.path.exists(file_path):
                        new_filename = f"{base_name}_{counter}{extension}"
                        file_path = os.path.join(upload_path, new_filename)
                        counter += 1
                    
                    # Save file
                    file.save(file_path)
                    uploaded_count += 1
                    logger.info(f"Uploaded file: {file_path}")
                
                else:
                    errors.append(f"File type not allowed: {file.filename}")
            
            except Exception as e:
                error_msg = f"Error uploading {file.filename}: {str(e)}"
                errors.append(error_msg)
                logger.error(error_msg)
        
        # Show results
        if uploaded_count > 0:
            flash(f'Successfully uploaded {uploaded_count} file(s)', 'success')
        
        if errors:
            for error in errors:
                flash(error, 'error')
        
        if uploaded_count == 0 and not errors:
            flash('No valid files were uploaded', 'error')
    
    except Exception as e:
        logger.error(f"Error in upload route: {str(e)}")
        flash(f'Upload error: {str(e)}', 'error')
    
    return redirect(url_for('index'))

@app.route('/category/<path:category_path>')
def show_category(category_path):
    """Show files in a specific category."""
    try:
        full_path = os.path.join(UPLOAD_FOLDER, category_path)
        
        if not os.path.exists(full_path) or not os.path.isdir(full_path):
            flash('Category not found', 'error')
            return redirect(url_for('index'))
        
        # Get breadcrumbs
        path_parts = category_path.split('/')
        breadcrumbs = []
        current_path = ""
        
        for part in path_parts:
            if current_path:
                current_path += "/" + part
            else:
                current_path = part
            breadcrumbs.append({'name': part, 'path': current_path})
        
        # Get files and subdirectories
        files = []
        subcategories = []
        
        for item in os.listdir(full_path):
            item_path = os.path.join(full_path, item)
            
            if os.path.isfile(item_path):
                file_info = {
                    'name': item,
                    'type': get_file_type(item),
                    'size': os.path.getsize(item_path),
                    'modified': os.path.getmtime(item_path)
                }
                files.append(file_info)
            
            elif os.path.isdir(item_path):
                subcategories.append(item)
        
        # Sort files by name
        files.sort(key=lambda x: x['name'].lower())
        subcategories.sort()
        
        category_name = path_parts[-1]
        
        return render_template('category.html',
                             category_name=category_name,
                             category_path=category_path,
                             breadcrumbs=breadcrumbs,
                             files=files,
                             subcategories=subcategories,
                             categories=get_categories())
    
    except Exception as e:
        logger.error(f"Error showing category {category_path}: {str(e)}")
        flash(f'Error loading category: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/download/<path:file_path>')
def download_file(file_path):
    """Download a specific file."""
    try:
        full_path = os.path.join(UPLOAD_FOLDER, file_path)
        
        if not os.path.exists(full_path) or not os.path.isfile(full_path):
            flash('File not found', 'error')
            return redirect(url_for('index'))
        
        return send_file(full_path, as_attachment=True)
    
    except Exception as e:
        logger.error(f"Error downloading file {file_path}: {str(e)}")
        flash(f'Download error: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/download_category/<path:category_path>')
def download_category(category_path):
    """Download entire category as ZIP."""
    try:
        full_path = os.path.join(UPLOAD_FOLDER, category_path)
        
        if not os.path.exists(full_path) or not os.path.isdir(full_path):
            flash('Category not found', 'error')
            return redirect(url_for('index'))
        
        # Create temporary ZIP file
        zip_filename = f"{category_path.replace('/', '_')}_archive.zip"
        zip_path = os.path.join(UPLOAD_FOLDER, zip_filename)
        
        with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for root, dirs, files in os.walk(full_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, full_path)
                    zipf.write(file_path, arcname)
        
        def remove_file(response):
            try:
                os.remove(zip_path)
            except Exception:
                pass
            return response
        
        return send_file(zip_path, as_attachment=True, 
                        download_name=zip_filename)
    
    except Exception as e:
        logger.error(f"Error creating ZIP for category {category_path}: {str(e)}")
        flash(f'Archive creation error: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/delete_category/<path:category_path>', methods=['POST'])
def delete_category(category_path):
    """Delete entire category."""
    try:
        full_path = os.path.join(UPLOAD_FOLDER, category_path)
        
        if not os.path.exists(full_path):
            flash('Category not found', 'error')
            return redirect(url_for('index'))
        
        import shutil
        shutil.rmtree(full_path)
        flash(f'Category "{category_path}" deleted successfully', 'success')
    
    except Exception as e:
        logger.error(f"Error deleting category {category_path}: {str(e)}")
        flash(f'Delete error: {str(e)}', 'error')
    
    return redirect(url_for('index'))

@app.route('/view_document/<path:file_path>')
def view_document(file_path):
    """View document content in browser."""
    try:
        full_path = os.path.join(UPLOAD_FOLDER, file_path)
        
        if not os.path.exists(full_path) or not os.path.isfile(full_path):
            flash('File not found', 'error')
            return redirect(url_for('index'))
        
        filename = os.path.basename(file_path)
        file_type = os.path.splitext(filename)[1].lower()
        
        # Extract text content
        content = extract_text_from_file(full_path)
        
        return render_template('document_viewer.html',
                             filename=filename,
                             file_path=file_path,
                             file_type=file_type,
                             content=content)
    
    except Exception as e:
        logger.error(f"Error viewing document {file_path}: {str(e)}")
        flash(f'Document viewing error: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/move_file', methods=['POST'])
def move_file():
    """Move file to different category."""
    try:
        source_path = request.form.get('source_path', '').strip()
        target_category = request.form.get('target_category', '').strip()
        target_subcategory = request.form.get('target_subcategory', '').strip()
        
        if not source_path or not target_category:
            flash('Missing required information for file move', 'error')
            return redirect(url_for('index'))
        
        # Sanitize paths
        target_category = secure_filename(target_category)
        if target_subcategory:
            target_subcategory = secure_filename(target_subcategory)
        
        source_full_path = os.path.join(UPLOAD_FOLDER, source_path)
        
        if not os.path.exists(source_full_path):
            flash('Source file not found', 'error')
            return redirect(url_for('index'))
        
        # Create target directory structure
        if target_subcategory:
            target_dir = os.path.join(UPLOAD_FOLDER, target_category, target_subcategory)
            target_path = os.path.join(target_category, target_subcategory, os.path.basename(source_path))
        else:
            target_dir = os.path.join(UPLOAD_FOLDER, target_category)
            target_path = os.path.join(target_category, os.path.basename(source_path))
        
        os.makedirs(target_dir, exist_ok=True)
        
        target_full_path = os.path.join(UPLOAD_FOLDER, target_path)
        
        # Handle duplicate filenames
        counter = 1
        base_name, extension = os.path.splitext(os.path.basename(source_path))
        
        while os.path.exists(target_full_path):
            new_filename = f"{base_name}_{counter}{extension}"
            if target_subcategory:
                target_path = os.path.join(target_category, target_subcategory, new_filename)
            else:
                target_path = os.path.join(target_category, new_filename)
            target_full_path = os.path.join(UPLOAD_FOLDER, target_path)
            counter += 1
        
        # Move the file
        import shutil
        shutil.move(source_full_path, target_full_path)
        
        # Clean up empty directories
        source_dir = os.path.dirname(source_full_path)
        try:
            if not os.listdir(source_dir):
                os.rmdir(source_dir)
        except:
            pass
        
        target_display = f"{target_category}"
        if target_subcategory:
            target_display += f"/{target_subcategory}"
        
        flash(f'File moved to {target_display} successfully', 'success')
        
        # Redirect to target category
        if target_subcategory:
            return redirect(url_for('show_category', category_path=f"{target_category}/{target_subcategory}"))
        else:
            return redirect(url_for('show_category', category_path=target_category))
    
    except Exception as e:
        logger.error(f"Error moving file {source_path}: {str(e)}")
        flash(f'Move error: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/copy_file', methods=['POST'])
def copy_file():
    """Copy file to different category."""
    try:
        source_path = request.form.get('source_path', '').strip()
        target_category = request.form.get('target_category', '').strip()
        target_subcategory = request.form.get('target_subcategory', '').strip()
        
        if not source_path or not target_category:
            flash('Missing required information for file copy', 'error')
            return redirect(url_for('index'))
        
        # Sanitize paths
        target_category = secure_filename(target_category)
        if target_subcategory:
            target_subcategory = secure_filename(target_subcategory)
        
        source_full_path = os.path.join(UPLOAD_FOLDER, source_path)
        
        if not os.path.exists(source_full_path):
            flash('Source file not found', 'error')
            return redirect(url_for('index'))
        
        # Create target directory structure
        if target_subcategory:
            target_dir = os.path.join(UPLOAD_FOLDER, target_category, target_subcategory)
            target_path = os.path.join(target_category, target_subcategory, os.path.basename(source_path))
        else:
            target_dir = os.path.join(UPLOAD_FOLDER, target_category)
            target_path = os.path.join(target_category, os.path.basename(source_path))
        
        os.makedirs(target_dir, exist_ok=True)
        
        target_full_path = os.path.join(UPLOAD_FOLDER, target_path)
        
        # Handle duplicate filenames
        counter = 1
        base_name, extension = os.path.splitext(os.path.basename(source_path))
        
        while os.path.exists(target_full_path):
            new_filename = f"{base_name}_copy_{counter}{extension}"
            if target_subcategory:
                target_path = os.path.join(target_category, target_subcategory, new_filename)
            else:
                target_path = os.path.join(target_category, new_filename)
            target_full_path = os.path.join(UPLOAD_FOLDER, target_path)
            counter += 1
        
        # Copy the file
        import shutil
        shutil.copy2(source_full_path, target_full_path)
        
        target_display = f"{target_category}"
        if target_subcategory:
            target_display += f"/{target_subcategory}"
        
        flash(f'File copied to {target_display} successfully', 'success')
        
        # Redirect to target category
        if target_subcategory:
            return redirect(url_for('show_category', category_path=f"{target_category}/{target_subcategory}"))
        else:
            return redirect(url_for('show_category', category_path=target_category))
    
    except Exception as e:
        logger.error(f"Error copying file {source_path}: {str(e)}")
        flash(f'Copy error: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/api/categories')
def api_categories():
    """API endpoint to get all categories."""
    try:
        categories = get_categories()
        category_list = list(categories.keys())
        return jsonify({'categories': category_list})
    except Exception as e:
        logger.error(f"Error getting categories via API: {str(e)}")
        return jsonify({'error': str(e)}), 500

@app.route('/create_category', methods=['POST'])
def create_category():
    """Create new category or subcategory."""
    try:
        category_name = request.form.get('category_name', '').strip()
        parent_category = request.form.get('parent_category', '').strip()
        
        if not category_name:
            flash('Category name is required', 'error')
            return redirect(url_for('index'))
        
        # Sanitize category name
        category_name = secure_filename(category_name)
        
        if parent_category:
            parent_category = secure_filename(parent_category)
            full_path = os.path.join(UPLOAD_FOLDER, parent_category, category_name)
            category_path = f"{parent_category}/{category_name}"
        else:
            full_path = os.path.join(UPLOAD_FOLDER, category_name)
            category_path = category_name
        
        if os.path.exists(full_path):
            flash('Category already exists', 'error')
            return redirect(url_for('index'))
        
        os.makedirs(full_path, exist_ok=True)
        flash(f'Category "{category_path}" created successfully', 'success')
        
        return redirect(url_for('show_category', category_path=category_path))
    
    except Exception as e:
        logger.error(f"Error creating category: {str(e)}")
        flash(f'Category creation error: {str(e)}', 'error')
        return redirect(url_for('index'))

@app.route('/delete/<path:file_path>', methods=['POST'])
def delete_file(file_path):
    """Delete a specific file."""
    try:
        full_path = os.path.join(UPLOAD_FOLDER, file_path)
        
        if not os.path.exists(full_path):
            flash('File not found', 'error')
            return redirect(url_for('index'))
        
        os.remove(full_path)
        flash(f'File "{os.path.basename(file_path)}" deleted successfully', 'success')
        
        # Clean up empty directories
        source_dir = os.path.dirname(full_path)
        try:
            if not os.listdir(source_dir):
                os.rmdir(source_dir)
        except:
            pass
        
        # Redirect to category if we know it
        category_path = '/'.join(file_path.split('/')[:-1])
        if category_path:
            return redirect(url_for('show_category', category_path=category_path))
        else:
            return redirect(url_for('index'))
    
    except Exception as e:
        logger.error(f"Error deleting file {file_path}: {str(e)}")
        flash(f'Delete error: {str(e)}', 'error')
        return redirect(url_for('index'))

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)
